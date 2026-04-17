"""PowerPoint parser using python-pptx."""

from __future__ import annotations

from pathlib import Path

from app.models.mcp import ParsedDocument


async def parse_pptx(file_path: str | Path) -> ParsedDocument:
    """Parse a PPTX file and extract slide text."""
    try:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text)
            if texts:
                slides_text.append(f"## Slide {i}\n" + "\n".join(texts))

        return ParsedDocument(
            source=str(file_path),
            content="\n\n".join(slides_text),
            file_type=".pptx",
            pages=len(prs.slides),
            metadata={"slide_count": len(prs.slides)},
        )
    except ImportError:
        return ParsedDocument(source=str(file_path), content="Error: python-pptx not installed", file_type=".pptx")
